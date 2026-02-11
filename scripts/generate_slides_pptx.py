#!/usr/bin/env python3
"""
Generate AstrID Midterm Presentation PPTX from slide content.

Usage:
    python scripts/generate_slides_pptx.py

Output:
    docs/research/MIDTERM_SLIDES.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Color Palette ──────────────────────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)        # Deep navy/charcoal
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)    # Bright sky blue
ACCENT_GOLD = RGBColor(0xFF, 0xD5, 0x4F)    # Gold/amber
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x99)
HIGHLIGHT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
HIGHLIGHT_RED = RGBColor(0xEF, 0x53, 0x50)
SUBTLE_BG = RGBColor(0x22, 0x22, 0x3A)      # Slightly lighter bg for cards


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, bullets, left, top, width, height,
                     font_size=17, color=WHITE, line_spacing=1.4,
                     bullet_color=None):
    """Add a bulleted text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (indented with leading spaces or tabs)
        indent_level = 0
        clean_bullet = bullet
        if bullet.startswith("  ") or bullet.startswith("\t"):
            indent_level = 1
            clean_bullet = bullet.strip()

        p.text = clean_bullet
        p.font.size = Pt(font_size - (2 * indent_level))
        p.font.color.rgb = color if indent_level == 0 else LIGHT_GRAY
        p.font.name = "Calibri"
        p.level = indent_level
        p.space_after = Pt(4)
        p.space_before = Pt(2)

    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_number_badge(slide, number, left, top):
    """Add a small circular slide number badge."""
    size = Inches(0.4)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def make_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(1), Inches(1.2), Inches(8), Inches(0.9),
                "AstrID", font_size=52, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(2.0), Inches(8), Inches(0.6),
                "Supernova Detection Pipeline", font_size=28, color=WHITE,
                bold=False, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.6), Inches(8), Inches(0.5),
                "Machine Learning for Astronomical Transient Identification",
                font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(3.5), Inches(3.2), Inches(3),
                    color=ACCENT_GOLD)

    # Author info
    info_lines = [
        "Chris Lawrence",
        "CSCI 491 — Senior Research Project",
        "Midterm Presentation — February 2026",
        "",
        "Advisors: Prof. Luis Meneses  •  Prof. Gregory Arkos"
    ]
    add_bullet_slide(slide, info_lines,
                     Inches(1), Inches(3.6), Inches(8), Inches(2.0),
                     font_size=15, color=LIGHT_GRAY)

    # Fix alignment for author block
    for p in slide.shapes[-1].text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER


def make_content_slide(prs, slide_num, title, bullets, notes="",
                       highlight_indices=None, two_column=False,
                       col1_bullets=None, col2_bullets=None,
                       col1_title=None, col2_title=None):
    """Create a standard content slide with title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Slide number badge
    add_number_badge(slide, slide_num, Inches(0.3), Inches(0.25))

    # Title
    add_textbox(slide, Inches(0.85), Inches(0.2), Inches(8.5), Inches(0.65),
                title, font_size=28, color=ACCENT_BLUE, bold=True)

    # Accent line
    add_accent_line(slide, Inches(0.85), Inches(0.85), Inches(2.5))

    if two_column and col1_bullets and col2_bullets:
        # Two-column layout
        if col1_title:
            add_textbox(slide, Inches(0.6), Inches(1.05), Inches(4.2), Inches(0.4),
                        col1_title, font_size=16, color=ACCENT_GOLD, bold=True)
        if col2_title:
            add_textbox(slide, Inches(5.2), Inches(1.05), Inches(4.2), Inches(0.4),
                        col2_title, font_size=16, color=ACCENT_GOLD, bold=True)

        y_start = Inches(1.45) if col1_title else Inches(1.1)
        add_bullet_slide(slide, col1_bullets,
                         Inches(0.6), y_start, Inches(4.2), Inches(4.0),
                         font_size=15, color=WHITE)
        add_bullet_slide(slide, col2_bullets,
                         Inches(5.2), y_start, Inches(4.2), Inches(4.0),
                         font_size=15, color=WHITE)
    else:
        # Standard single column
        add_bullet_slide(slide, bullets,
                         Inches(0.85), Inches(1.1), Inches(8.3), Inches(4.5),
                         font_size=17, color=WHITE)

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def make_section_slide(prs, title, subtitle=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, SUBTLE_BG)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.0),
                title, font_size=40, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.5),
                    subtitle, font_size=18, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Decorative lines
    add_accent_line(slide, Inches(3), Inches(2.3), Inches(4), ACCENT_GOLD)
    add_accent_line(slide, Inches(3), Inches(4.2), Inches(4), ACCENT_GOLD)


def make_questions_slide(prs):
    """Final Questions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.0),
                "Questions?", font_size=48, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(3.5), Inches(3.1), Inches(3), ACCENT_GOLD)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.4),
                "Thank you!", font_size=22, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.4),
                "Chris Lawrence  •  CSCI 491  •  February 2026",
                font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


def build_presentation():
    """Build the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    make_title_slide(prs)

    # ── Slide 2: Agenda ──
    make_content_slide(prs, 2, "What We'll Cover Today", [
        "The problem: scale of modern astronomical surveys",
        "AstrID's pipeline architecture (5 stages)",
        "Key discovery: same-mission requirement",
        "Image differencing technique and results",
        "Current dataset and proof-of-concept triplets",
        "Lessons learned and next steps",
    ], notes="Give the audience a map of the talk. ~20 min total. Save questions for the end or ask as we go.")

    # ── Section: The Problem ──
    make_section_slide(prs, "The Problem", "Why does astronomy need automation?")

    # ── Slide 3: Scale of Modern Astronomy ──
    make_content_slide(prs, 3, "Why Automation Matters", [
        "Modern surveys (e.g., Rubin/LSST) will generate ~10 million alerts per night",
        "Far too much data for human inspection",
        "Transient events (supernovae, novae) appear briefly and fade",
        "Missing a supernova = losing a unique scientific opportunity",
        "Automation is not a convenience — it is a necessity",
    ], notes="Emphasize sheer scale. LSST alone will produce 10M alerts/night. Humans can't keep up.")

    # ── Slide 4: What Are Supernovae ──
    make_content_slide(prs, 4, "Supernovae: Exploding Stars", [
        "Catastrophic explosions at the end of a star's life",
        "Briefly outshine their entire host galaxy (billions of stars)",
        "Critical for cosmology: Type Ia SNe are 'standard candles' for measuring cosmic distances",
        "Extremely rare: ~1 per century per galaxy",
        "Detection requires comparing images taken at different times",
    ], notes="Supernovae are among the most energetic events in the universe — scientifically invaluable but rare and fleeting.")

    # ── Slide 5: The Bogus Problem ──
    make_content_slide(prs, 5, "The False Positive Challenge", [
        "Classical detection flags 73–629 candidates per image at 5σ threshold",
        "Only ~1 is a real supernova",
        "The rest: cosmic rays, subtraction artifacts, variable stars, galaxy nuclei",
        "A model predicting 'always bogus' gets 99.99% accuracy",
        "… but is scientifically useless",
        "Need ML to separate real transients from artifacts",
    ], notes="This is the core ML motivation. Classical algorithms can't distinguish real from bogus alone.")

    # ── Section: AstrID Pipeline ──
    make_section_slide(prs, "AstrID Pipeline", "Architecture and design")

    # ── Slide 6: What is AstrID ──
    make_content_slide(prs, 6, "AstrID: Project Overview", [
        "End-to-end pipeline for automated supernova detection",
        "Combines classical astronomy techniques with machine learning",
        "Full workflow: data acquisition → preprocessing → comparison → detection → classification",
        "Built from scratch — not using existing survey pipelines",
        "Emphasis on hands-on experimentation and system design",
    ], notes="AstrID handles the entire problem. Building it from scratch gives deep understanding of each component.")

    # ── Slide 7: Pipeline Architecture ──
    make_content_slide(prs, 7, "AstrID Pipeline: 5 Stages", [
        "Stage 1 — Query: Retrieve observation metadata from MAST archive",
        "Stage 2 — Filter: Ensure same-mission compatibility",
        "Stage 3 — Download: Acquire FITS files (60–80% size reduction via smart filtering)",
        "Stage 4 — Organize: Create training-ready directory structures",
        "Stage 5 — Differencing: Full astronomical image differencing pipeline",
        "",
        "Output → difference images + significance maps + detection masks",
    ], notes="Walk through each stage. Modular design lets us test and debug each independently.")

    # ── Slide 8: Data Acquisition ──
    make_content_slide(prs, 8, "Data Acquisition: The Supernova Catalog", [
        "Compiled catalog of 6,542 known supernovae with coordinates + discovery dates",
        "Query MAST (Mikulski Archive for Space Telescopes)",
        "Download reference images (before SN) and science images (after discovery)",
        "Missions: SWIFT, GALEX, PS1 (Pan-STARRS)",
        "Temporal windowing ensures correct before/after separation",
    ], notes="We start with known supernovae for ground truth. MAST is NASA's main archive for space telescope data.")

    # ── Slide 9: FITS Files ──
    make_content_slide(prs, 9, "Working with FITS Files", [
        "FITS (Flexible Image Transport System) — standard astronomical image format",
        "Contains pixel data + rich metadata (WCS coordinates, exposure, filters)",
        "Complex: extensions, HDUs, 3D arrays, compressed variants (.fits.gz)",
        "WCS maps pixels to sky coordinates (RA/Dec)",
        "Discovering compressed FITS support → 296% increase in training pairs (56 → 222)",
    ], notes="FITS files are complex. Discovering compressed files nearly tripled our dataset — a major breakthrough.")

    # ── Section: Key Discovery ──
    make_section_slide(prs, "Critical Discovery", "The same-mission requirement")

    # ── Slide 10: Pilot Study ──
    make_content_slide(prs, 10, "Pilot Study: 19 Supernovae", [
        "Started with 19 known supernovae from the catalog",
        "Downloaded reference + science images from MAST",
        "Attempted image differencing on all pairs",
        "Only 8 of 19 (42%) had usable same-mission pairs",
        "11 had cross-mission data (e.g., PS1 reference + SWIFT science)",
        "Cross-mission differencing produced massive, unusable artifacts",
    ], notes="The pilot study was crucial — more than half the data was unusable. This led to our most important discovery.")

    # ── Slide 11: Same-Mission Requirement ──
    make_content_slide(prs, 11, "Why Cross-Mission Differencing Fails", [
        "Different telescopes have fundamentally different characteristics:",
        "  • Point Spread Functions (PSFs)",
        "  • Filter bandpasses and sensitivities",
        "  • Pixel scales and field of view",
        "  • Noise properties and detector artifacts",
        "Subtracting cross-mission images → massive artifacts swamp real signals",
        "Requirement: same mission AND same filter for every pair",
    ], notes="Single most important discovery. Obvious in hindsight but fundamentally reshaped the pipeline. You can't compare apples and oranges.")

    # ── Slide 12: Same-Mission Impact ──
    make_content_slide(prs, 12, "Impact of the Same-Mission Constraint", [
        "Initial approach: download everything, filter later → 75% wasted",
        "New approach: filter at query stage before downloading",
        "Download size reduced 60–80%",
        "Success rate: 25% → 75–100% for complete training pairs",
        "Lesson: understanding data constraints early saves enormous effort",
    ], notes="Once we understood this, we built it into the pipeline. Huge efficiency gain.")

    # ── Section: Image Differencing ──
    make_section_slide(prs, "Image Differencing", "The core detection technique")

    # ── Slide 13: What Is Differencing ──
    make_content_slide(prs, 13, "What Is Image Differencing?", [
        "Classical astronomical technique:",
        "  Subtract a reference image from a science image",
        "Anything that changed between epochs appears as residual signal",
        "Simple pixel subtraction is inadequate — massive artifacts",
        "Must carefully: align, match PSFs, normalize flux, subtract backgrounds",
        "Done correctly → only real changes remain",
    ], notes="Conceptually simple — subtract two images. But without careful preprocessing, artifacts dominate.")

    # ── Slide 14: 9-Stage Pipeline ──
    make_content_slide(prs, 14, "Differencing Pipeline: 9 Stages", [
        "1. FITS File Loading (robust HDU handling)",
        "2. WCS Alignment (sub-pixel reprojection)",
        "3. Background Estimation & Subtraction",
        "4. PSF Estimation (FWHM from bright stars)",
        "5. PSF Matching (Gaussian convolution)",
        "6. Flux Normalization (robust median matching)",
        "7. Difference Image Computation",
        "8. Significance Map Generation (SNR per pixel)",
        "9. Source Detection (DAOStarFinder, 5σ threshold)",
    ], notes="Nine stages, each addressing a specific challenge. Skip any one and the result is garbage.")

    # ── Slide 15: WCS Alignment ──
    make_content_slide(prs, 15, "WCS Alignment: Sub-Pixel Precision", [
        "Images taken at different times have different pointings and rotations",
        "Even 1-pixel misalignment → massive dipole residuals around every source",
        "Solution: reproject science image onto reference WCS grid",
        "Uses reproject.reproject_interp() for sub-pixel accuracy",
        "Result: 93.6%–100% overlap on processed pairs",
    ], notes="Arguably the most critical step. Sub-pixel accuracy is essential — even tiny offsets create bright artifacts.")

    # ── Slide 16: PSF Matching ──
    make_content_slide(prs, 16, "PSF Matching & Background Subtraction", [
        "Background subtraction: remove varying sky levels",
        "  Uses photutils.Background2D with median estimator (50–64 px boxes)",
        "PSF estimation: measure FWHM from bright stars",
        "  Typically 3.5–4.0 px for SWIFT UVOT",
        "PSF matching: convolve sharper image with Gaussian kernel",
        "Without matching: bright stars leave large residuals",
        "With matching: residuals dramatically reduced → transients revealed",
    ], notes="The PSF is the telescope's fingerprint. Different observations have different seeing. Match them so stars cancel cleanly.")

    # ── Slide 17: Significance Maps ──
    make_content_slide(prs, 17, "Significance Maps & Source Detection", [
        "Significance map: each pixel = signal-to-noise ratio",
        "  Value = difference / combined_noise",
        "Enables consistent thresholding across images with different noise",
        "'Whitened noise' — uniform statistical properties",
        "DAOStarFinder on significance map at 5σ threshold",
        "Shape filters (sharpness bounds) reject cosmic rays + extended sources",
        "Output: candidate transient positions with significance values",
    ], notes="Instead of raw pixel thresholds, every pixel tells us how many σ above noise. 5σ = 1-in-3.5 million chance of random noise.")

    # ── Section: Results ──
    make_section_slide(prs, "Results", "What has the pipeline found?")

    # ── Slide 18: SN 2014J ──
    make_content_slide(prs, 18, "Detection: SN 2014J (2120σ)", [
        "SN 2014J: brightest supernova in decades (galaxy M82)",
        "Detected at 2120σ significance — unmistakable signal",
        "SWIFT UVOT U-band filter",
        "Found at pixel (447.9, 555.4) — matches known position",
        "73 total candidates in image; 1 is the real supernova",
        "Demonstrates: pipeline works correctly when inputs are good",
    ], notes="Our best example. 2120 times above noise. But even here, 72 other false detections.")

    # ── Slide 19: Batch Results ──
    make_content_slide(prs, 19, "Batch Processing: 5 Supernovae", [
        "5 SNe fully processed: 2014J, 2014ai, 2014bh, 2014bi, 2014cs",
        "All SWIFT UVOT U-band (uuu filter)",
        "Max significance range: 412σ to 2120σ",
        "Detection counts: 73–629 candidates per image",
        "Image overlap: 66%–100% after WCS alignment",
        "Consistent pipeline behavior across multiple targets",
    ], notes="Consistent results across multiple targets. Wide range in detection counts shows how variable the false positive problem is.")

    # ── Slide 20: Production Dataset ──
    make_content_slide(prs, 20, "Production Dataset Statistics", [
        "222 complete supernova pairs (reference + science)",
        "2,282 FITS files total (1,170 reference, 1,112 science)",
        "Mission breakdown:",
        "  GALEX: 105 pairs",
        "  SWIFT: 17+ pairs",
        "  PS1: additional pairs",
        "99.1% download success rate",
        "15 YAML configs for reproducible generation",
    ], notes="Solid foundation. 222 pairs gives enough to start training. High success rate proves pipeline reliability.")

    # ── Slide 21: False Positive Problem ──
    make_content_slide(prs, 21, "The False Positive Problem", [
        "5σ threshold: 73–629 candidates, only ~1 real per image",
        "False positive sources:",
        "  • Cosmic rays (sharp single-pixel spikes)",
        "  • Subtraction artifacts (imperfect alignment)",
        "  • Variable stars (real variability, not SNe)",
        "  • Host galaxy nuclei (AGN activity)",
        "Accuracy is the WRONG metric",
        "Correct: precision, recall, F1, AUCPR",
    ], notes="This bridges to ML. Classical detection does physics perfectly but can't tell SN from cosmic ray. That's the CNN's job.")

    # ── Section: ML Preparation ──
    make_section_slide(prs, "ML Preparation", "From detection to classification")

    # ── Slide 22: Training Triplets ──
    make_content_slide(prs, 22, "Image Triplets for CNN Input", [
        "Triplet = 63×63 pixel cutout of (science, reference, difference)",
        "3-channel input — analogous to RGB for a standard CNN",
        "Labels: Real (1) = known SN position; Bogus (0) = artifacts",
        "Proof of concept: 57 samples (30 real, 27 bogus)",
        "Augmentation: rotation, flipping",
        "Pipeline built and tested — needs to run at scale",
    ], notes="This bridges classical detection and ML. Small stamps around each detection. 57 samples prove it works; need thousands for real training.")

    # ── Slide 23: Config System ──
    make_content_slide(prs, 23, "Configuration-Driven Reproducibility", [
        "All pipeline runs driven by YAML configuration files",
        "15 configs for different missions, filters, temporal windows",
        "Two execution modes:",
        "  run_pipeline_from_config.py — full 5-stage batch",
        "  run_pipeline_per_sn.py — one SN at a time, checkpoint/resume",
        "Any result can be exactly reproduced with the same config",
        "Scientific reproducibility built into the architecture",
    ], notes="Reproducibility is central. Every dataset can be recreated. Modular design lets us test stages independently.")

    # ── Slide 24: Training Strategy ──
    make_content_slide(prs, 24, "Mission-Specific Training Strategy", [
        "Config system enables sophisticated experiments:",
        "  1. Train separate models per mission (SWIFT, GALEX, PS1)",
        "  2. Train combined model for generalization",
        "  3. Cross-mission validation: train SWIFT+PS1, test GALEX",
        "  4. Filter-specific models for specialized wavelengths",
        "Goal: learn transient physics, not instrument artifacts",
        "Domain shift is a fundamental challenge in astronomical ML",
    ], notes="Modular configs enable real scientific experiments. Train on one telescope, test on another — proves the model learned physics, not artifacts.")

    # ── Slide 25: Literature Alignment ──
    make_content_slide(prs, 25, "Alignment with Published Research",
        bullets=[],
        two_column=True,
        col1_title="Implemented (Phases 1–4)",
        col1_bullets=[
            "Data acquisition from MAST archive",
            "WCS alignment & calibration",
            "Image differencing with PSF matching",
            "Source detection (DAOStarFinder 5σ)",
            "Triplet generation (proof of concept)",
        ],
        col2_title="Remaining (Phases 5–8)",
        col2_bullets=[
            "CNN classifier (Braai-style)",
            "Training with class imbalance handling",
            "Evaluation (AUCPR, precision-recall)",
            "Deployment & real-time inference",
            "Active learning loop (stretch)",
        ],
        notes="Following ZTF/Braai methodology. Phases 1-4 done, 5-8 are second half focus.")

    # ── Section: Lessons & Next Steps ──
    make_section_slide(prs, "Lessons & Next Steps", "What we learned and where we're going")

    # ── Slide 26: Lessons Learned ──
    make_content_slide(prs, 26, "Key Lessons from the First Half", [
        "Data quality > model complexity — better preprocessing beat fancier models",
        "Discover constraints early — same-mission requirement reshaped everything",
        "Small details matter — compressed FITS → 296% more training pairs",
        "Invest in infrastructure — modular design + configs accelerated all later work",
        "Research is non-linear — progress came from revisiting assumptions",
        "Debugging ML ≠ debugging software — emergent behavior from data-parameter interactions",
    ], notes="These lessons are as important as results. Biggest gains came from understanding data better, not fancier algorithms.")

    # ── Slide 27: Next Steps - Scale ──
    make_content_slide(prs, 27, "Next: Scale the Dataset", [
        "Current: 5 fully processed SNe, 57 triplet samples",
        "Target: 200+ difference image sets with quality filtering",
        "Enforce minimum 85% overlap threshold",
        "Run differencing across all 222 pairs",
        "Expand SN coordinate catalog for better labeling",
        "Deliverable: large-scale difference images + processing summaries",
    ], notes="Immediate priority: scaling up. We have the pipeline — now run it at scale. 200+ difference images for meaningful ML training.")

    # ── Slide 28: Next Steps - CNN ──
    make_content_slide(prs, 28, "Next: CNN Classifier",
        bullets=[],
        two_column=True,
        col1_title="Phase B: Triplets at Scale",
        col1_bullets=[
            "Generate thousands of 63×63 triplets",
            "Balanced: ~50% real / 50% bogus",
            "Augmentation: rotation, flip",
            "Save in NPZ format",
            "Train/val/test splits",
        ],
        col2_title="Phase C: Braai-Style CNN",
        col2_bullets=[
            "Input: 63×63×3 (sci, ref, diff)",
            "Conv2D(32)→(64)→(128)",
            "Dense(256)→Sigmoid",
            "Output: P(real) ∈ [0, 1]",
            "TensorFlow/Keras or PyTorch",
        ],
        notes="Generate triplets, then train a CNN. Architecture is well-established. Challenge is quality training data.")

    # ── Slide 29: Next Steps - Training ──
    make_content_slide(prs, 29, "Next: Training & Evaluation", [
        "Temporal/spatial splits to avoid data leakage",
        "Class imbalance handling: weighted loss or oversampling",
        "Metrics: precision, recall, F1, AUCPR — not accuracy",
        "Threshold tuning: minimize false negatives (missing a SN is catastrophic)",
        "Stretch: cross-mission validation, real-time API, active learning",
    ], notes="Evaluation is where science meets engineering. Right metrics, right splits, right threshold.")

    # ── Slide 30: Timeline ──
    make_content_slide(prs, 30, "Remaining Timeline", [
        "Feb–Mar 2026: Scale differencing to 200+ pairs; generate triplets",
        "Mar–Apr 2026: Implement and train CNN classifier; evaluate with AUCPR",
        "Apr 2026: Cross-mission validation experiments",
        "Apr–May 2026: (Stretch) Inference API + active learning prototype",
        "May 2026: Final report and presentation",
        "",
        "Foundation is solid — ML phase is the focus going forward",
    ], notes="Clear roadmap. Hardest part (understanding data, building pipeline) is done. Remaining work is primarily ML.")

    # ── Slide 31: Conclusion ──
    make_content_slide(prs, 31, "Summary & Key Takeaways", [
        "Built a production-ready 5-stage pipeline from scratch",
        "Discovered critical constraints through experimentation",
        "Implemented 9-stage differencing pipeline — detections at 412σ–2120σ",
        "Created 222 training pairs (2,282 FITS) across 3 missions",
        "Proved triplet generation at small scale (57 samples)",
        "Pipeline is verified, modular, and ready for ML scaling",
        "",
        "Key insight: data quality and infrastructure > model complexity",
    ], notes="First half established a solid foundation. Second half: scaling up and applying ML to the false positive problem.")

    # ── Slide 32: Questions ──
    make_questions_slide(prs)

    return prs


def main():
    import os
    prs = build_presentation()
    out_path = os.path.join(
        os.path.dirname(__file__), "..", "docs", "research", "MIDTERM_SLIDES.pptx"
    )
    out_path = os.path.abspath(out_path)
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
